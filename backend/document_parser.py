import os
import PyPDF2
try:
    from pptx import Presentation
except ImportError:
    pass

def parse_document(file_path: str) -> str:
    """Extrae el texto de archivos PDF, TXT, MD, o Powerpoint."""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
        elif ext in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif ext in ['.pptx', '.ppt']:
            try:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except Exception as e:
                print(f"Error procesando pptx: {e}. Asegurate de tener python-pptx instalado.")
        else:
            raise ValueError("Formato no soportado en el parser.")
            
    except Exception as e:
        print(f"Error parseando documento: {e}")
        return ""
        
    return text.strip()
